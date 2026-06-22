from __future__ import annotations

import argparse
import csv
import json
import math
from collections import Counter, defaultdict
from datetime import datetime
from pathlib import Path
from typing import Any

import matplotlib.pyplot as plt
from PIL import Image
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_AUTO_SIZE
from pptx.util import Inches, Pt


ROOT = Path(__file__).resolve().parents[1]
DEFAULT_OUTPUT = ROOT / "reports" / "MoonRockStack_academic_report_20260618.pptx"


COLORS = {
    "ink": RGBColor(30, 35, 38),
    "muted": RGBColor(92, 99, 105),
    "line": RGBColor(210, 214, 218),
    "bg": RGBColor(247, 248, 250),
    "dark": RGBColor(22, 26, 29),
    "moon": RGBColor(189, 173, 124),
    "earth": RGBColor(71, 124, 167),
    "accent": RGBColor(154, 96, 52),
    "danger": RGBColor(170, 73, 62),
    "ok": RGBColor(64, 132, 92),
}


FONT = "Microsoft YaHei"


def main() -> int:
    args = parse_args()
    output = args.output.resolve()
    output.parent.mkdir(parents=True, exist_ok=True)
    stats = collect_stats(ROOT)
    chart_dir = output.parent / "ppt_assets"
    chart_dir.mkdir(parents=True, exist_ok=True)
    charts = make_charts(stats, chart_dir)
    ppt = build_presentation(stats, charts)
    ppt.save(output)
    summary_path = output.with_suffix(".summary.json")
    summary_path.write_text(json.dumps(stats["summary"], indent=2, ensure_ascii=False) + "\n", encoding="utf-8")
    print(output)
    print(summary_path)
    return 0


def parse_args() -> argparse.Namespace:
    parser = argparse.ArgumentParser(description="Create the MoonStack academic progress report PPT.")
    parser.add_argument("--output", type=Path, default=DEFAULT_OUTPUT)
    return parser.parse_args()


def collect_stats(root: Path) -> dict[str, Any]:
    batch = root / "batch_runs"
    results = read_all_csv(batch, "results.csv")
    placements = read_all_csv(batch, "placement_log.csv")
    candidate_poses = read_all_csv(batch, "candidate_pose_log.csv")
    failures = read_all_csv(batch, "failure_cases.csv")
    captures = read_all_csv(batch, "capture_manifest.csv")

    for row in results:
        row["_height"] = to_float(row.get("stack_height_m"))
        row["_success"] = to_int(row.get("success"))
        row["_shape_success"] = to_int(row.get("shape_success"))
        row["_stable_count"] = to_float(row.get("stable_count"))
        row["_rock_count"] = to_float(row.get("rock_count"))
        row["_stable_fraction"] = row["_stable_count"] / max(row["_rock_count"], 1.0)
        row["_rmse"] = to_float(row.get("target_rmse_xy_m"))
        row["_drift"] = to_float(row.get("max_horizontal_drift_m"))
        row["_velocity"] = to_float(row.get("velocity_inf_norm"))

    total_runs = len(results)
    strict_success = sum(row["_success"] for row in results)
    shape_success = sum(row["_shape_success"] for row in results)
    run_dirs = {row.get("_dir", "") for row in results}
    image_paths = []
    for row in captures:
        image_paths.extend(path for path in str(row.get("images", "")).split() if path)

    by_gravity = aggregate_key(results, "gravity")
    by_target = aggregate_key(results, "target_name")
    by_strategy = aggregate_key(results, "strategy")
    failure_counter = Counter()
    failure_target_counter = Counter()
    for row in failures:
        reason = row.get("failure_reason", "") or "unlabeled_failure"
        failure_counter[reason] += 1
        failure_target_counter[(row.get("target_name", "unknown"), reason)] += 1

    key_runs = {
        "moon_4course_partial": find_result(
            results,
            "20260618_single_face_wall_4course_v1_assignment_plan_wall_bonded_gated_smoke",
            gravity="moon",
        ),
        "earth_4course_success": find_result(
            results,
            "20260618_single_face_wall_4course_v1_assignment_fallback5_gated_smoke",
            gravity="earth",
        ),
        "moon_4course_fallback_failure": find_result(
            results,
            "20260618_single_face_wall_4course_v1_assignment_fallback5_gated_smoke",
            gravity="moon",
        ),
        "earth_high_fast": find_result(
            results,
            "20260618_high_single_face_wall_v1_earth_fast",
            gravity="earth",
        ),
        "moon_high_fast": find_result(
            results,
            "20260618_high_single_face_wall_v1_moon_fast",
            gravity="moon",
        ),
        "earth_ranker_high": find_result(
            results,
            "20260618_high_single_face_wall_v1_earth_ranker_top3",
            gravity="earth",
        ),
        "moon_neural_high": find_result(
            results,
            "20260618_neural_stonefit_pose_high_wall_top8_earth_moon_v1",
            gravity="moon",
        ),
        "earth_neural_high": find_result(
            results,
            "20260618_neural_stonefit_pose_high_wall_top8_earth_moon_v1",
            gravity="earth",
        ),
        "moon_neural_column": find_result(
            results,
            "20260618_neural_stonefit_pose_multi_column_lit_earth_moon_v1",
            gravity="moon",
        ),
        "earth_neural_column": find_result(
            results,
            "20260618_neural_stonefit_pose_multi_column_lit_earth_moon_v1",
            gravity="earth",
        ),
    }

    model_metrics = read_json(
        root / "batch_runs" / "20260618_modular_small_networks_v2_high_and_ranker" / "metrics.json"
    )
    dataset_summary = read_json(
        root / "batch_runs" / "20260618_learning_dataset_v3_high_and_ranker" / "dataset_summary.json"
    )
    ranker_eval = read_json(
        root / "batch_runs" / "20260618_candidate_pose_ranker_eval_v2_high" / "summary.json"
    )

    image_map = {
        "fourcourse_earth_front": first_existing(
            [
                root
                / "batch_runs"
                / "20260618_single_face_wall_4course_v1_assignment_fallback5_gated_smoke"
                / "captures_960x720"
                / "00_single_face_wall_4course_v1_success_wall_bonded_earth_trial_00"
                / "wall_front_rgb.png",
                *sorted(
                    (
                        root
                        / "batch_runs"
                        / "20260618_single_face_wall_4course_v1_assignment_fallback5_gated_smoke"
                        / "captures_960x720"
                    ).glob("*earth*/*front_rgb.png")
                ),
            ]
        ),
        "fourcourse_earth_depth": first_existing(
            sorted(
                (
                    root
                    / "batch_runs"
                    / "20260618_single_face_wall_4course_v1_assignment_fallback5_gated_smoke"
                    / "captures_960x720"
                ).glob("*earth*/*wall_top_depth.png")
            )
        ),
        "moon_partial_front": first_existing(
            sorted(
                (
                    root
                    / "batch_runs"
                    / "20260618_single_face_wall_4course_v1_assignment_plan_wall_bonded_gated_smoke"
                    / "captures_960x720"
                ).glob("*moon*/*wall_front_rgb.png")
            )
        ),
        "neural_high_moon_front": root
        / "batch_runs"
        / "20260618_neural_stonefit_pose_high_wall_top8_earth_moon_v1"
        / "captures_960x720"
        / "01_single_face_wall_high_v1_failure_statics_wall_moon_trial_00"
        / "wall_front_rgb.png",
        "neural_high_moon_top_depth": root
        / "batch_runs"
        / "20260618_neural_stonefit_pose_high_wall_top8_earth_moon_v1"
        / "captures_960x720"
        / "01_single_face_wall_high_v1_failure_statics_wall_moon_trial_00"
        / "wall_top_depth.png",
        "neural_column_moon_front": root
        / "batch_runs"
        / "20260618_neural_stonefit_pose_multi_column_lit_earth_moon_v1"
        / "captures_960x720"
        / "01_multi_stone_column_v3_failure_literature_column_moon_trial_00"
        / "front_rgb.png",
        "neural_column_moon_top_depth": root
        / "batch_runs"
        / "20260618_neural_stonefit_pose_multi_column_lit_earth_moon_v1"
        / "captures_960x720"
        / "01_multi_stone_column_v3_failure_literature_column_moon_trial_00"
        / "top_depth.png",
    }

    summary = {
        "created_at": datetime.now().isoformat(timespec="seconds"),
        "result_rows": total_runs,
        "run_dirs_with_results": len(run_dirs),
        "strict_success": strict_success,
        "shape_success": shape_success,
        "strict_success_rate": strict_success / max(total_runs, 1),
        "shape_success_rate": shape_success / max(total_runs, 1),
        "placement_log_rows": len(placements),
        "candidate_pose_rows": len(candidate_poses),
        "failure_case_rows": len(failures),
        "capture_cases": len(captures),
        "capture_images": len(image_paths),
        "papers_unique_core": 11,
        "ppt_master_status": "GitHub direct and mirror clone failed; report generated with python-pptx.",
    }

    return {
        "results": results,
        "placements": placements,
        "candidate_poses": candidate_poses,
        "failures": failures,
        "captures": captures,
        "by_gravity": by_gravity,
        "by_target": by_target,
        "by_strategy": by_strategy,
        "failure_counter": failure_counter,
        "failure_target_counter": failure_target_counter,
        "key_runs": key_runs,
        "model_metrics": model_metrics,
        "dataset_summary": dataset_summary,
        "ranker_eval": ranker_eval,
        "image_map": image_map,
        "summary": summary,
    }


def read_all_csv(batch: Path, filename: str) -> list[dict[str, str]]:
    rows: list[dict[str, str]] = []
    for path in sorted(batch.rglob(filename)):
        try:
            with path.open("r", encoding="utf-8", newline="") as handle:
                for row in csv.DictReader(handle):
                    row["_path"] = str(path)
                    row["_dir"] = str(path.parent)
                    row["_run_name"] = path.parent.name
                    rows.append(row)
        except UnicodeDecodeError:
            with path.open("r", encoding="utf-8-sig", newline="") as handle:
                for row in csv.DictReader(handle):
                    row["_path"] = str(path)
                    row["_dir"] = str(path.parent)
                    row["_run_name"] = path.parent.name
                    rows.append(row)
    return rows


def read_json(path: Path) -> dict[str, Any]:
    if not path.exists():
        return {}
    return json.loads(path.read_text(encoding="utf-8"))


def aggregate_key(results: list[dict[str, Any]], key: str) -> list[dict[str, Any]]:
    grouped: dict[str, list[dict[str, Any]]] = defaultdict(list)
    for row in results:
        grouped[str(row.get(key, "") or "unknown")].append(row)
    output = []
    for value, rows in grouped.items():
        count = len(rows)
        success = sum(row["_success"] for row in rows)
        shape = sum(row["_shape_success"] for row in rows)
        output.append(
            {
                key: value,
                "count": count,
                "success": success,
                "shape_success": shape,
                "success_rate": success / max(count, 1),
                "shape_success_rate": shape / max(count, 1),
                "max_height": max((row["_height"] for row in rows), default=0.0),
                "mean_height": sum(row["_height"] for row in rows) / max(count, 1),
                "mean_stable_fraction": sum(row["_stable_fraction"] for row in rows) / max(count, 1),
            }
        )
    return sorted(output, key=lambda row: (-row["count"], str(row[key])))


def find_result(results: list[dict[str, Any]], run_name: str, gravity: str | None = None) -> dict[str, Any]:
    matches = [row for row in results if row.get("_run_name") == run_name]
    if gravity is not None:
        matches = [row for row in matches if row.get("gravity") == gravity]
    if not matches:
        return {}
    return sorted(matches, key=lambda row: (-row["_height"], row["_rmse"]))[0]


def first_existing(paths: list[Path]) -> Path | None:
    for path in paths:
        if path and path.exists():
            return path
    return None


def to_float(value: Any, default: float = 0.0) -> float:
    try:
        if value in {"", None}:
            return default
        out = float(value)
        if math.isnan(out) or math.isinf(out):
            return default
        return out
    except (TypeError, ValueError):
        return default


def to_int(value: Any, default: int = 0) -> int:
    try:
        if value in {"", None}:
            return default
        return int(float(value))
    except (TypeError, ValueError):
        return default


def make_charts(stats: dict[str, Any], chart_dir: Path) -> dict[str, Path]:
    charts: dict[str, Path] = {}
    plt.style.use("seaborn-v0_8-whitegrid")

    summary = stats["summary"]
    labels = ["Strict success", "Shape success", "Strict failure"]
    values = [
        summary["strict_success"],
        summary["shape_success"],
        summary["result_rows"] - summary["strict_success"],
    ]
    path = chart_dir / "overall_success.png"
    fig, ax = plt.subplots(figsize=(7.2, 3.8), dpi=180)
    ax.bar(labels, values, color=["#40845c", "#b29a57", "#aa493e"])
    ax.set_ylabel("Runs")
    ax.set_title("Overall structured-run outcomes")
    for i, value in enumerate(values):
        ax.text(i, value + max(values) * 0.02, str(value), ha="center", va="bottom", fontsize=10)
    fig.tight_layout()
    fig.savefig(path, transparent=False)
    plt.close(fig)
    charts["overall_success"] = path

    targets = stats["by_target"][:10]
    path = chart_dir / "target_counts.png"
    fig, ax = plt.subplots(figsize=(8.8, 4.6), dpi=180)
    y = list(range(len(targets)))
    ax.barh(y, [row["count"] for row in targets], color="#477ca7", label="runs")
    ax.barh(y, [row["success"] for row in targets], color="#40845c", label="strict success")
    ax.set_yticks(y)
    ax.set_yticklabels([short_label(row["target_name"], 28) for row in targets], fontsize=8)
    ax.invert_yaxis()
    ax.set_xlabel("Runs")
    ax.set_title("Top targets by number of trials")
    ax.legend(loc="lower right")
    fig.tight_layout()
    fig.savefig(path, transparent=False)
    plt.close(fig)
    charts["target_counts"] = path

    reasons = stats["failure_counter"].most_common(9)
    path = chart_dir / "failure_reasons.png"
    fig, ax = plt.subplots(figsize=(8.8, 4.4), dpi=180)
    y = list(range(len(reasons)))
    ax.barh(y, [count for _reason, count in reasons], color="#aa493e")
    ax.set_yticks(y)
    ax.set_yticklabels([short_label(reason, 32) for reason, _count in reasons], fontsize=8)
    ax.invert_yaxis()
    ax.set_xlabel("Failure rows")
    ax.set_title("Dominant logged failure modes")
    fig.tight_layout()
    fig.savefig(path, transparent=False)
    plt.close(fig)
    charts["failure_reasons"] = path

    key = stats["key_runs"]
    height_items = [
        ("Moon 4-course partial", key["moon_4course_partial"]),
        ("Earth 4-course success", key["earth_4course_success"]),
        ("Earth high wall", key["earth_high_fast"]),
        ("Moon high wall", key["moon_high_fast"]),
        ("Moon neural high wall", key["moon_neural_high"]),
        ("Moon neural column", key["moon_neural_column"]),
    ]
    height_items = [(name, row) for name, row in height_items if row]
    path = chart_dir / "milestone_heights.png"
    fig, ax = plt.subplots(figsize=(8.8, 4.2), dpi=180)
    ax.bar([name for name, _row in height_items], [row["_height"] for _name, row in height_items], color="#b29a57")
    ax.set_ylabel("Height (m)")
    ax.set_title("Milestone height comparison")
    ax.tick_params(axis="x", rotation=30, labelsize=8)
    for i, (_name, row) in enumerate(height_items):
        ax.text(i, row["_height"] + 0.01, f"{row['_height']:.3f}", ha="center", fontsize=8)
    fig.tight_layout()
    fig.savefig(path, transparent=False)
    plt.close(fig)
    charts["milestone_heights"] = path

    return charts


def short_label(text: str, limit: int) -> str:
    if len(text) <= limit:
        return text
    return text[: limit - 1] + "..."


def build_presentation(stats: dict[str, Any], charts: dict[str, Path]) -> Presentation:
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    add_title_slide(prs, stats)
    add_problem_slide(prs)
    add_literature_slide(prs)
    add_priors_slide(prs)
    add_geometry_slide(prs)
    add_environment_slide(prs)
    add_pipeline_slide(prs)
    add_data_overview_slide(prs, stats, charts)
    add_phase4_slide(prs, stats)
    add_fourcourse_images_slide(prs, stats)
    add_high_wall_slide(prs, stats)
    add_column_slide(prs, stats)
    add_neural_method_slide(prs)
    add_model_metrics_slide(prs, stats)
    add_neural_wall_slide(prs, stats)
    add_neural_column_slide(prs, stats)
    add_failure_slide(prs, stats, charts)
    add_experience_slide(prs)
    add_limitations_slide(prs)
    add_next_steps_slide(prs)
    add_appendix_slide(prs, stats)
    return prs


def blank_slide(prs: Presentation, title: str, kicker: str | None = None):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(255, 255, 255)
    add_header(slide, title, kicker)
    return slide


def add_header(slide, title: str, kicker: str | None = None) -> None:
    if kicker:
        box = slide.shapes.add_textbox(Inches(0.55), Inches(0.22), Inches(12.2), Inches(0.22))
        p = box.text_frame.paragraphs[0]
        p.text = kicker
        p.font.name = FONT
        p.font.size = Pt(8.5)
        p.font.bold = True
        p.font.color.rgb = COLORS["accent"]
    box = slide.shapes.add_textbox(Inches(0.55), Inches(0.45), Inches(12.2), Inches(0.55))
    p = box.text_frame.paragraphs[0]
    p.text = title
    p.font.name = FONT
    p.font.size = Pt(25)
    p.font.bold = True
    p.font.color.rgb = COLORS["ink"]
    line = slide.shapes.add_shape(1, Inches(0.55), Inches(1.08), Inches(12.2), Inches(0.02))
    line.fill.solid()
    line.fill.fore_color.rgb = COLORS["line"]
    line.line.fill.background()


def add_title_slide(prs: Presentation, stats: dict[str, Any]) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = COLORS["dark"]
    add_text(
        slide,
        "月面干式石头堆叠实验阶段汇报",
        0.72,
        0.65,
        9.7,
        0.68,
        size=31,
        color=RGBColor(245, 244, 238),
        bold=True,
    )
    add_text(
        slide,
        "从分形多面体石头生成、文献先验、MuJoCo 地/月重力对照，到模块化小网络选石/选位",
        0.75,
        1.42,
        10.9,
        0.45,
        size=15,
        color=RGBColor(219, 215, 201),
    )
    img = stats["image_map"].get("neural_high_moon_front")
    if img and img.exists():
        add_image(slide, img, 6.85, 2.05, 5.65, 4.15)
    summary = stats["summary"]
    metric_cards(
        slide,
        [
            ("结构实验", f"{summary['result_rows']} 次"),
            ("放置记录", f"{summary['placement_log_rows']} 条"),
            ("候选位姿", f"{summary['candidate_pose_rows']} 条"),
            ("核心资料", f"{summary['papers_unique_core']} 篇/部"),
        ],
        x=0.75,
        y=2.25,
        card_w=1.35,
        card_h=0.74,
        dark=True,
    )
    add_text(slide, "2026-06-18 | MoonStack / MuJoCo / Lunar landmark stacking", 0.75, 6.78, 9.5, 0.25, size=10, color=RGBColor(190, 190, 184))


def add_problem_slide(prs: Presentation) -> None:
    slide = blank_slide(prs, "研究目标：不是石头堆，而是可解释结构", "Problem")
    add_section_box(
        slide,
        0.72,
        1.35,
        5.7,
        4.8,
        "短期目标",
        [
            "月面路标堆叠：形成可识别、可定位的小型干式石结构。",
            "当前主目标：单面墙、高墙、单/多石柱。",
            "判定不只看高度，还看目标占位、漂移、速度和稳定石数。",
        ],
    )
    add_section_box(
        slide,
        6.78,
        1.35,
        5.7,
        4.8,
        "长期目标",
        [
            "月面原位资源利用：不依赖砂浆或标准构件。",
            "从路标扩展到干砌墙、围护结构，最终服务月面建筑。",
            "核心科学问题：低重力怎样改变不规则石块接触、沉降和结构强度。",
        ],
    )
    add_text(slide, "汇报立场：失败案例是数据，不是噪声；当前所有严格失败都会保留并进入学习数据。", 0.9, 6.45, 11.7, 0.35, size=13, color=COLORS["danger"], bold=True)


def add_literature_slide(prs: Presentation) -> None:
    slide = blank_slide(prs, "统计和吸收的核心文献/资料", "Literature")
    columns = ["资料", "方法重点", "转化为本实验的先验"]
    rows = [
        ["Furrer 2017", "online next-best object/pose", "每块石头多姿态搜索 + 物理验证"],
        ["Johns 2020", "真实干石墙在线规划", "边感知、边规划、边修正"],
        ["Liu 2018", "DQN dry stacking", "限制动作空间后再学习策略"],
        ["Liu 2021", "仿真筛选 + 真实堆叠", "稳定姿态库、顺序规划"],
        ["Liu & Napp 2023", "稳定性序列规划", "不能只优化单步，要评估整体强度"],
        ["Menezes 2021", "model-free RL from depth", "可由观测学习目标体积/稳定性"],
        ["Zhang 2020 / Wang 2024", "抓取-堆叠联合学习", "抓取成功要服务后续堆叠"],
        ["GraspNet v2 / Efficient-WAM", "数据集与世界模型", "下一阶段采集数据训练小网络/世界模型"],
        ["Hibbeler Statics", "静力学、受力、摩擦", "支撑多边形、重心、摩擦角、倾覆风险"],
    ]
    add_table(slide, columns, rows, 0.55, 1.28, 12.25, 5.55, font_size=8.5)
    add_text(slide, "本阶段使用 11 篇/部核心资料：干式堆叠论文、抓取/强化学习论文、GraspNet/WAM 方向和静力学教材。", 0.72, 6.92, 11.8, 0.25, size=9.5, color=COLORS["muted"])


def add_priors_slide(prs: Presentation) -> None:
    slide = blank_slide(prs, "从文献到可执行堆叠先验", "Priors")
    headers = ["石头先验", "结构先验", "搜索/学习先验"]
    items = [
        [
            "角砾/多面体、宽 fracture faces",
            "拒绝尖刺、极扁 slab、光滑球",
            "按体积、轴比、compactness、spike_score 聚类",
        ],
        [
            "底层：大、宽、低尖刺、高 compactness",
            "中层：破缝、搭接、支撑重心覆盖",
            "墙：batter/收分、tie stones；柱：环形/芯石支撑",
        ],
        [
            "先生成有限候选，再 MuJoCo settle",
            "顺序规划要优化后续稳定性",
            "网络先做排序/剪枝，不能过早替代物理验证",
        ],
    ]
    x_positions = [0.72, 4.62, 8.52]
    for x, header, bullet_items in zip(x_positions, headers, items):
        add_section_box(slide, x, 1.35, 3.45, 4.95, header, bullet_items)


def add_geometry_slide(prs: Presentation) -> None:
    slide = blank_slide(prs, "石头生成与几何分类：多凸/多面体，不要尖刺", "Geometry")
    add_section_box(
        slide,
        0.65,
        1.28,
        5.9,
        5.2,
        "当前生成类别",
        [
            "equant / subangular / wedge / fractured / elongated clast",
            "后续扩展：bearing_block、course_block、tie_bridge、cap_block、interlock_block",
            "聚类不是展示用途，而是进入 base/middle/tie/cap 的角色选择。",
        ],
    )
    add_section_box(
        slide,
        6.75,
        1.28,
        5.9,
        5.2,
        "筛选阈值",
        [
            "spike_score > 0.16: reject",
            "flatness > 1.62 或 short_to_mid < 0.62: reject",
            "elongation > 1.85 / compactness < 0.22: reject",
            "四层墙 catalog: 360 生成，337 接受，23 尖刺拒绝，0 slab 拒绝。",
        ],
    )


def add_environment_slide(prs: Presentation) -> None:
    slide = blank_slide(prs, "仿真环境与物理校准", "Environment")
    metric_cards(
        slide,
        [
            ("Earth g", "9.80665 m/s²"),
            ("Moon g", "1.624 m/s²"),
            ("Friction", "mu=1.15"),
            ("Angle", "48.99°"),
        ],
        x=0.65,
        y=1.35,
        card_w=2.55,
        card_h=0.9,
    )
    add_section_box(
        slide,
        0.75,
        2.72,
        5.7,
        3.65,
        "验证项目",
        [
            "自由落体拟合：地球/月球加速度与配置一致。",
            "倾斜等效摩擦：阈值接近 atan(1.15)。",
            "同一批石头、同一目标、同一随机种子下比较地/月重力。",
        ],
    )
    add_section_box(
        slide,
        6.75,
        2.72,
        5.7,
        3.65,
        "限制",
        [
            "MuJoCo 是刚体接触近似。",
            "未建模月壤 cohesion、粉尘、碎裂、颗粒嵌锁和接触老化。",
            "当前结论应定位为刚体接触趋势和策略对比。",
        ],
    )


def add_pipeline_slide(prs: Presentation) -> None:
    slide = blank_slide(prs, "自动实验 Pipeline", "Method")
    steps = [
        "生成角砾石库",
        "提取几何特征",
        "聚类/角色评分",
        "目标结构 slots",
        "候选石头/姿态",
        "MuJoCo settle",
        "稳定性/失败日志",
        "数据集/小网络",
    ]
    x0, y = 0.55, 2.35
    w, h, gap = 1.35, 0.72, 0.18
    for i, step in enumerate(steps):
        x = x0 + i * (w + gap)
        shape = slide.shapes.add_shape(5, Inches(x), Inches(y), Inches(w), Inches(h))
        shape.fill.solid()
        shape.fill.fore_color.rgb = RGBColor(238, 241, 244) if i < 4 else RGBColor(235, 227, 204)
        shape.line.color.rgb = COLORS["line"]
        tf = shape.text_frame
        tf.clear()
        p = tf.paragraphs[0]
        p.text = step
        p.font.name = FONT
        p.font.size = Pt(10)
        p.font.bold = True
        p.alignment = PP_ALIGN.CENTER
        if i < len(steps) - 1:
            add_text(slide, "→", x + w + 0.035, y + 0.2, 0.1, 0.2, size=16, color=COLORS["muted"], bold=True)
    add_section_box(
        slide,
        0.8,
        4.1,
        11.7,
        1.75,
        "关键记录项",
        [
            "results.csv: 每个结构实验的 success、shape_success、height、RMSE、drift、velocity。",
            "placement_log.csv / failure_cases.csv: 每块石头的角色、落点、支撑、失败原因。",
            "candidate_pose_log.csv: 位姿候选、ranker 概率、是否被搜索选中。",
            "captures_960x720: RGB、深度 PNG、深度 NPY、正视和俯视图。",
        ],
    )


def add_data_overview_slide(prs: Presentation, stats: dict[str, Any], charts: dict[str, Path]) -> None:
    slide = blank_slide(prs, "总体实验规模与成功率", "Data")
    summary = stats["summary"]
    metric_cards(
        slide,
        [
            ("结构实验", f"{summary['result_rows']}"),
            ("严格成功", f"{summary['strict_success']} ({summary['strict_success_rate']:.1%})"),
            ("形状成功", f"{summary['shape_success']} ({summary['shape_success_rate']:.1%})"),
            ("失败记录", f"{summary['failure_case_rows']}"),
        ],
        x=0.65,
        y=1.25,
        card_w=2.85,
        card_h=0.9,
    )
    add_image(slide, charts["overall_success"], 0.75, 2.52, 5.65, 3.35)
    add_image(slide, charts["target_counts"], 6.65, 2.2, 5.95, 3.95)
    add_text(
        slide,
        f"额外数据：placement_log={summary['placement_log_rows']} 条，candidate_pose_log={summary['candidate_pose_rows']} 条，capture_cases={summary['capture_cases']}，图像/深度文件约 {summary['capture_images']} 个。总体成功率包含早期简单/烟测任务；高墙和多石柱阶段仍为 strict failure。",
        0.75,
        6.62,
        11.7,
        0.32,
        size=10,
        color=COLORS["muted"],
    )


def add_phase4_slide(prs: Presentation, stats: dict[str, Any]) -> None:
    slide = blank_slide(prs, "阶段性成果 I：四层单面墙", "Milestone")
    rows = []
    for label, key in [
        ("Moon clean partial", "moon_4course_partial"),
        ("Earth gated success", "earth_4course_success"),
        ("Moon fallback failure", "moon_4course_fallback_failure"),
    ]:
        row = stats["key_runs"].get(key, {})
        rows.append(result_table_row(label, row))
    add_table(
        slide,
        ["Case", "Placed", "Stable", "Courses", "Height", "RMSE", "Drift", "Success"],
        rows,
        0.65,
        1.28,
        12.0,
        1.4,
        font_size=8.5,
    )
    add_section_box(
        slide,
        0.75,
        3.05,
        5.75,
        3.55,
        "阶段意义",
        [
            "已经能堆出 target-locked 四层结构，而不是散乱土堆。",
            "Moon clean partial: 14/24 放置、14/14 稳定、RMSE 0.045 m、drift 0.0015 m。",
            "Earth fallback 首次达到 current-gate success，但仍是短段墙。",
        ],
    )
    add_section_box(
        slide,
        6.75,
        3.05,
        5.75,
        3.55,
        "暴露问题",
        [
            "Moon 下 Earth fallback 不直接迁移，横向漂移和 side outlier 增加。",
            "跳过 slot 说明缺少 support repair。",
            "真实任务不能靠大量随机尝试，需要学习化候选排序。",
        ],
    )


def add_fourcourse_images_slide(prs: Presentation, stats: dict[str, Any]) -> None:
    slide = blank_slide(prs, "四层墙图像：正视与俯视深度", "Evidence")
    add_labeled_image(slide, stats["image_map"].get("fourcourse_earth_front"), 0.65, 1.28, 5.85, 4.35, "Earth gated success / wall-front RGB")
    add_labeled_image(slide, stats["image_map"].get("fourcourse_earth_depth"), 6.82, 1.28, 5.85, 4.35, "Earth gated success / wall-top depth")
    add_text(slide, "注：四层墙中 Earth 有 current-gate success；Moon 更干净的线形结果仍是 partial，因为跳过了 10 个 slots。", 0.78, 6.25, 11.5, 0.42, size=11, color=COLORS["muted"])


def add_high_wall_slide(prs: Presentation, stats: dict[str, Any]) -> None:
    slide = blank_slide(prs, "阶段性成果 II：更高单面墙", "High wall")
    rows = []
    for label, key in [
        ("Earth fast", "earth_high_fast"),
        ("Moon fast", "moon_high_fast"),
        ("Earth pose-ranker top3", "earth_ranker_high"),
        ("Moon neural top8", "moon_neural_high"),
    ]:
        rows.append(result_table_row(label, stats["key_runs"].get(key, {})))
    add_table(
        slide,
        ["Case", "Placed", "Stable", "Courses", "Height", "RMSE", "Drift", "Success"],
        rows,
        0.65,
        1.28,
        12.0,
        1.78,
        font_size=8.3,
    )
    add_section_box(
        slide,
        0.75,
        3.38,
        11.75,
        2.75,
        "结论",
        [
            "8 个 visible courses 已经可以稳定触达，但 strict success 仍为 0。",
            "Ranker 可以推高高度，但若只模仿局部 pose search，会恶化 RMSE/drift。",
            "Moon neural high wall 放满 31/31，height=0.283 m，drift=0.141 m；失败主因是残余速度和 outlier。",
        ],
    )


def add_column_slide(prs: Presentation, stats: dict[str, Any]) -> None:
    slide = blank_slide(prs, "阶段性成果 III：单石柱与多石柱", "Column")
    rows = [
        ["Single-column v4 Earth", "10", "2 / 10", "10", "0.237 m", "0.280 m", "0.012 m", "0"],
        ["Single-column v4 Moon", "10", "2 / 10", "10", "0.212 m", "0.232 m", "0.233 m", "0"],
        result_table_row("Multi-column neural Earth", stats["key_runs"].get("earth_neural_column", {})),
        result_table_row("Multi-column neural Moon", stats["key_runs"].get("moon_neural_column", {})),
    ]
    add_table(
        slide,
        ["Case", "Placed", "Stable", "Courses", "Height", "RMSE", "Drift", "Success"],
        rows,
        0.65,
        1.28,
        12.0,
        1.78,
        font_size=8.3,
    )
    add_section_box(
        slide,
        0.75,
        3.38,
        11.75,
        2.75,
        "结论",
        [
            "一块石头一层的单柱非常脆弱：可见 10 层，但稳定石头只有 2/10。",
            "多石柱/环形+芯石是更合理的高度路线。",
            "Moon neural multi-column height=0.303 m，40/40 放置、23/40 稳定，但 target spread 太大。",
        ],
    )


def add_neural_method_slide(prs: Presentation) -> None:
    slide = blank_slide(prs, "神经网络化：先做模块化小网络，而不是一步端到端", "Learning")
    add_section_box(
        slide,
        0.65,
        1.25,
        3.85,
        4.95,
        "StoneFitNet",
        [
            "输入：石头几何、目标 role、gravity、target slot。",
            "输出：这块石头是否适合当前 slot。",
            "在线用途：石头池排序，不能过窄剪枝。",
        ],
    )
    add_section_box(
        slide,
        4.8,
        1.25,
        3.85,
        4.95,
        "CandidatePoseRankNet",
        [
            "输入：候选位置/四元数、石头特征、slot 信息。",
            "输出：是否接近当前搜索赢家。",
            "在线用途：pose top-K 预筛，再交给 MuJoCo。",
        ],
    )
    add_section_box(
        slide,
        8.95,
        1.25,
        3.35,
        4.95,
        "MuJoCo 保底",
        [
            "网络只排序/剪枝。",
            "最终仍用物理 settle + hold 验证。",
            "下一步加入延迟稳定性标签。",
        ],
    )


def add_model_metrics_slide(prs: Presentation, stats: dict[str, Any]) -> None:
    slide = blank_slide(prs, "模块化小网络：训练数据与指标", "Model metrics")
    summary = stats["dataset_summary"]
    models = stats["model_metrics"].get("models", {})
    metric_cards(
        slide,
        [
            ("Run examples", str(summary.get("run_example_count", 480))),
            ("Placements", str(summary.get("placement_example_count", 8080))),
            ("Candidate poses", str(summary.get("candidate_pose_example_count", 1193))),
            ("Assignments", str(summary.get("assignment_candidate_example_count", 264))),
        ],
        x=0.65,
        y=1.25,
        card_w=2.85,
        card_h=0.86,
    )
    rows = []
    for name in ["stone_fit_net", "pose_accept_net", "moon_drift_risk_net", "candidate_pose_rank_net"]:
        item = models.get(name, {})
        rows.append(
            [
                name,
                str(item.get("row_count", "")),
                f"{to_float(item.get('accuracy')):.3f}" if "accuracy" in item else "",
                f"{to_float(item.get('precision')):.3f}" if "precision" in item else "",
                f"{to_float(item.get('recall')):.3f}" if "recall" in item else "",
                f"{to_float(item.get('f1')):.3f}" if "f1" in item else "",
            ]
        )
    add_table(slide, ["Model", "Rows", "Acc", "Precision", "Recall", "F1"], rows, 0.75, 2.55, 5.9, 2.45, font_size=8.8)
    ranker = stats["ranker_eval"]
    add_section_box(
        slide,
        7.0,
        2.4,
        5.35,
        2.78,
        "CandidatePoseRankNet Top-K",
        [
            f"候选位姿 rows: {ranker.get('candidate_pose_rows', 1193)}",
            f"候选组: {ranker.get('groups', 476)}",
            f"top-1 hit: {to_float(ranker.get('net_top1_hit_rate')):.3f}",
            f"top-3 hit: {to_float(ranker.get('net_top3_hit_rate')):.3f}",
            "注意：当前标签仍是模仿现有 hand-coded scorer。",
        ],
    )
    add_text(slide, "WorldDeltaNet 当前误差仍偏大，因为缺少局部支撑/高度图等 stack-state 特征。", 0.82, 6.25, 11.45, 0.35, size=11, color=COLORS["danger"], bold=True)


def add_neural_wall_slide(prs: Presentation, stats: dict[str, Any]) -> None:
    slide = blank_slide(prs, "神经化高墙：Moon 31/31 放置，但仍不是成功墙", "Neural wall")
    add_labeled_image(slide, stats["image_map"].get("neural_high_moon_front"), 0.65, 1.25, 6.15, 4.45, "Moon neural high wall / wall-front RGB")
    add_labeled_image(slide, stats["image_map"].get("neural_high_moon_top_depth"), 7.05, 1.25, 5.55, 4.45, "Moon neural high wall / wall-top depth")
    row = stats["key_runs"].get("moon_neural_high", {})
    text = (
        f"Moon neural high wall: placed={row.get('rock_count', '31')}/31, "
        f"stable={row.get('stable_count', '21')}/{row.get('rock_count', '31')}, "
        f"height={to_float(row.get('stack_height_m')):.3f} m, "
        f"RMSE={to_float(row.get('target_rmse_xy_m')):.3f} m, "
        f"drift={to_float(row.get('max_horizontal_drift_m')):.3f} m, "
        f"velocity={to_float(row.get('velocity_inf_norm')):.3f}; strict success=0."
    )
    add_text(slide, text, 0.78, 6.08, 11.8, 0.42, size=10.5, color=COLORS["ink"], bold=True)


def add_neural_column_slide(prs: Presentation, stats: dict[str, Any]) -> None:
    slide = blank_slide(prs, "神经化多石柱：高度更高，但形状太散", "Neural column")
    add_labeled_image(slide, stats["image_map"].get("neural_column_moon_front"), 0.65, 1.25, 6.15, 4.45, "Moon neural multi-stone column / RGB")
    add_labeled_image(slide, stats["image_map"].get("neural_column_moon_top_depth"), 7.05, 1.25, 5.55, 4.45, "Moon neural multi-stone column / top depth")
    row = stats["key_runs"].get("moon_neural_column", {})
    text = (
        f"Moon neural column: placed={row.get('rock_count', '40')}/40, "
        f"stable={row.get('stable_count', '23')}/{row.get('rock_count', '40')}, "
        f"height={to_float(row.get('stack_height_m')):.3f} m, "
        f"RMSE={to_float(row.get('target_rmse_xy_m')):.3f} m, "
        f"drift={to_float(row.get('max_horizontal_drift_m')):.3f} m; "
        "说明多石环/芯石有高度潜力，但需要 ring closure 和 compactness 目标。"
    )
    add_text(slide, text, 0.78, 6.08, 11.8, 0.42, size=10.5, color=COLORS["ink"], bold=True)


def add_failure_slide(prs: Presentation, stats: dict[str, Any], charts: dict[str, Path]) -> None:
    slide = blank_slide(prs, "失败模式统计：当前真正的瓶颈", "Failures")
    add_image(slide, charts["failure_reasons"], 0.65, 1.22, 6.25, 4.75)
    top = stats["failure_counter"].most_common(6)
    rows = [[reason, str(count)] for reason, count in top]
    add_table(slide, ["Failure reason", "Rows"], rows, 7.15, 1.42, 5.0, 2.2, font_size=8.5)
    add_section_box(
        slide,
        7.15,
        3.98,
        5.0,
        2.15,
        "解释",
        [
            "missed_target: 上层 outlier 或墙/柱形状散开。",
            "post_hold_drift: 初始看似可行，后续沉降/扰动后失稳。",
            "unstable_structure: 高度有了，但速度/漂移未收敛。",
        ],
    )


def add_experience_slide(prs: Presentation) -> None:
    slide = blank_slide(prs, "已形成的堆叠经验", "Lessons")
    add_section_box(
        slide,
        0.7,
        1.25,
        5.8,
        5.25,
        "什么石头放哪里",
        [
            "Base: 宽底、compact、低 spike、质量较大，优先 wedge/subangular/bearing blocks。",
            "Middle: 中等体积、低 flatness、能桥接下层 joints。",
            "Tie/core: 可适度 elongated，但不能 needle-like；需要跨 depth 或闭合 ring。",
            "Cap: 轻一些、compact、低残余速度，不追求单块最大高度。",
        ],
    )
    add_section_box(
        slide,
        6.8,
        1.25,
        5.8,
        5.25,
        "什么方法更有效",
        [
            "assignment plan 必须配 feasibility gate；盲目执行会制造 outlier。",
            "Moon 不能直接用 Earth fallback，需要更强 lateral drift gate。",
            "一石一层柱不可靠；多石环 + 芯石更适合高度。",
            "StoneFitNet 宽排序 + PoseRanker top-K + MuJoCo 验证，是当前最实用路线。",
        ],
    )


def add_limitations_slide(prs: Presentation) -> None:
    slide = blank_slide(prs, "科学限制与当前不可宣称内容", "Limitations")
    add_section_box(
        slide,
        0.75,
        1.28,
        11.6,
        5.25,
        "需要严谨表述",
        [
            "不能说已经能稳定堆高墙；当前高墙/高柱 strict success 仍为 0。",
            "当前网络标签主要来自已有启发式/MuJoCo 搜索，不是独立物理真值。",
            "刚体接触不等同于真实月壤接触；摩擦、粉尘、碎裂需要后续标定。",
            "图片中能看到部分结果仍接近局部堆散，因此必须用 target RMSE、wall spread、velocity 严格判别。",
            "阶段性成果是：结构化目标、数据闭环、可复现实验和在线小网络剪枝已经建立。",
        ],
    )


def add_next_steps_slide(prs: Presentation) -> None:
    slide = blank_slide(prs, "下一步：从局部搜索模仿转向延迟稳定性学习", "Next")
    add_section_box(
        slide,
        0.7,
        1.25,
        5.8,
        5.35,
        "数据采集",
        [
            "按 placement step 保存 RGB-D、局部高度图、支撑 centroid/span。",
            "记录 N 块石头之后是否仍稳定，形成 delayed stability label。",
            "扩大 Moon 高层负样本，专门训练 MoonDriftRiskNet。",
        ],
    )
    add_section_box(
        slide,
        6.8,
        1.25,
        5.8,
        5.35,
        "模型与策略",
        [
            "Course-closure classifier: 判断一层是否闭合/可继续加高。",
            "WorldDeltaNet 加入 stack-state，而不是只看单块几何。",
            "候选动作由网络提出，但 MuJoCo 继续做最终裁决。",
            "长期接入抓取：优化 grasp-place-stack 的联合目标。",
        ],
    )


def add_appendix_slide(prs: Presentation, stats: dict[str, Any]) -> None:
    slide = blank_slide(prs, "附：可复现输出与数据位置", "Appendix")
    summary = stats["summary"]
    add_section_box(
        slide,
        0.75,
        1.25,
        11.7,
        5.35,
        "本 PPT 使用的数据源",
        [
            r"D:\MoonStack\experiments\moon_rock_stack\docs\*.md",
            r"D:\MoonStack\experiments\moon_rock_stack\batch_runs\**\results.csv",
            r"D:\MoonStack\experiments\moon_rock_stack\batch_runs\**\placement_log.csv",
            r"D:\MoonStack\experiments\moon_rock_stack\batch_runs\**\candidate_pose_log.csv",
            r"D:\MoonStack\experiments\moon_rock_stack\batch_runs\**\captures_960x720",
            f"ppt-master 获取状态：{summary['ppt_master_status']}",
        ],
    )


def result_table_row(label: str, row: dict[str, Any]) -> list[str]:
    if not row:
        return [label, "", "", "", "", "", "", ""]
    placed = row.get("rock_count", "")
    stable = row.get("stable_count", "")
    return [
        label,
        str(placed),
        f"{stable} / {placed}",
        str(row.get("visible_courses", "")),
        f"{to_float(row.get('stack_height_m')):.3f} m",
        f"{to_float(row.get('target_rmse_xy_m')):.3f} m",
        f"{to_float(row.get('max_horizontal_drift_m')):.3f} m",
        str(row.get("success", "")),
    ]


def add_text(
    slide,
    text: str,
    x: float,
    y: float,
    w: float,
    h: float,
    size: float = 11,
    color: RGBColor | None = None,
    bold: bool = False,
    align=PP_ALIGN.LEFT,
) -> None:
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.clear()
    tf.margin_left = Inches(0.02)
    tf.margin_right = Inches(0.02)
    tf.margin_top = Inches(0.02)
    tf.margin_bottom = Inches(0.02)
    p = tf.paragraphs[0]
    p.text = text
    p.alignment = align
    p.font.name = FONT
    p.font.size = Pt(size)
    p.font.bold = bold
    p.font.color.rgb = color or COLORS["ink"]


def add_section_box(slide, x: float, y: float, w: float, h: float, title: str, bullets: list[str]) -> None:
    shape = slide.shapes.add_shape(1, Inches(x), Inches(y), Inches(w), Inches(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = COLORS["bg"]
    shape.line.color.rgb = COLORS["line"]
    add_text(slide, title, x + 0.18, y + 0.15, w - 0.35, 0.35, size=13, bold=True, color=COLORS["ink"])
    box = slide.shapes.add_textbox(Inches(x + 0.24), Inches(y + 0.62), Inches(w - 0.48), Inches(h - 0.8))
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
    for i, item in enumerate(bullets):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = item
        p.level = 0
        p.font.name = FONT
        p.font.size = Pt(10.5)
        p.font.color.rgb = COLORS["ink"]
        p.space_after = Pt(5)


def metric_cards(slide, cards: list[tuple[str, str]], x: float, y: float, card_w: float, card_h: float, dark: bool = False) -> None:
    for i, (label, value) in enumerate(cards):
        cx = x + i * (card_w + 0.22)
        shape = slide.shapes.add_shape(1, Inches(cx), Inches(y), Inches(card_w), Inches(card_h))
        shape.fill.solid()
        shape.fill.fore_color.rgb = RGBColor(44, 49, 52) if dark else COLORS["bg"]
        shape.line.color.rgb = RGBColor(80, 85, 88) if dark else COLORS["line"]
        add_text(slide, label, cx + 0.12, y + 0.12, card_w - 0.24, 0.22, size=8.5, color=RGBColor(190, 190, 184) if dark else COLORS["muted"], bold=True, align=PP_ALIGN.CENTER)
        add_text(slide, value, cx + 0.12, y + 0.39, card_w - 0.24, 0.28, size=13.5, color=RGBColor(245, 244, 238) if dark else COLORS["ink"], bold=True, align=PP_ALIGN.CENTER)


def add_table(slide, columns: list[str], rows: list[list[str]], x: float, y: float, w: float, h: float, font_size: float = 8.8) -> None:
    table_shape = slide.shapes.add_table(len(rows) + 1, len(columns), Inches(x), Inches(y), Inches(w), Inches(h))
    table = table_shape.table
    for col_idx, name in enumerate(columns):
        cell = table.cell(0, col_idx)
        cell.text = name
        set_cell_style(cell, RGBColor(226, 231, 235), bold=True, font_size=font_size)
    for row_idx, row in enumerate(rows, start=1):
        for col_idx, value in enumerate(row):
            cell = table.cell(row_idx, col_idx)
            cell.text = value
            fill = RGBColor(250, 251, 252) if row_idx % 2 == 1 else RGBColor(241, 244, 246)
            set_cell_style(cell, fill, bold=False, font_size=font_size)


def set_cell_style(cell, fill_color: RGBColor, bold: bool, font_size: float) -> None:
    cell.fill.solid()
    cell.fill.fore_color.rgb = fill_color
    cell.margin_left = Inches(0.04)
    cell.margin_right = Inches(0.04)
    cell.margin_top = Inches(0.02)
    cell.margin_bottom = Inches(0.02)
    for paragraph in cell.text_frame.paragraphs:
        paragraph.font.name = FONT
        paragraph.font.size = Pt(font_size)
        paragraph.font.bold = bold
        paragraph.font.color.rgb = COLORS["ink"]


def add_image(slide, path: Path, x: float, y: float, w: float, h: float) -> None:
    if not path or not path.exists():
        placeholder = slide.shapes.add_shape(1, Inches(x), Inches(y), Inches(w), Inches(h))
        placeholder.fill.solid()
        placeholder.fill.fore_color.rgb = RGBColor(235, 235, 235)
        placeholder.line.color.rgb = COLORS["line"]
        add_text(slide, "Image missing", x + 0.2, y + h / 2 - 0.15, w - 0.4, 0.3, size=12, color=COLORS["muted"], align=PP_ALIGN.CENTER)
        return
    with Image.open(path) as im:
        iw, ih = im.size
    box_ratio = w / h
    image_ratio = iw / ih
    if image_ratio >= box_ratio:
        shown_w = w
        shown_h = w / image_ratio
    else:
        shown_h = h
        shown_w = h * image_ratio
    px = x + (w - shown_w) / 2
    py = y + (h - shown_h) / 2
    slide.shapes.add_picture(str(path), Inches(px), Inches(py), width=Inches(shown_w), height=Inches(shown_h))
    border = slide.shapes.add_shape(1, Inches(x), Inches(y), Inches(w), Inches(h))
    border.fill.background()
    border.line.color.rgb = COLORS["line"]


def add_labeled_image(slide, path: Path | None, x: float, y: float, w: float, h: float, label: str) -> None:
    add_image(slide, path, x, y, w, h)
    add_text(slide, label, x, y + h + 0.08, w, 0.26, size=9, color=COLORS["muted"], bold=True, align=PP_ALIGN.CENTER)


if __name__ == "__main__":
    raise SystemExit(main())
