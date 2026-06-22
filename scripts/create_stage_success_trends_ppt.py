# -*- coding: utf-8 -*-
"""Create an enhanced group-meeting deck for the MoonStack wall work.

The script keeps the existing 2026-06-19 deck intact and appends a detailed
2026-06-22 update section focused on stage successes, failure-to-success data,
rule utility, neural-network integration, and layer-wise difficulties.
"""

from __future__ import annotations

import json
from pathlib import Path
from typing import Iterable

import matplotlib.pyplot as plt
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.util import Inches, Pt


ROOT = Path(r"D:\MoonStack")
BASE_PPTX = ROOT / "Asset" / "Reports" / "MoonRockStack_wall_group_meeting_20260619.pptx"
OUT_DIR = ROOT / "Asset" / "Reports"
ASSET_DIR = OUT_DIR / "ppt_assets_20260622_stage_success_trends"
OUT_PPTX = OUT_DIR / "MoonRockStack_wall_group_meeting_20260622_stage_success_trends.pptx"
OUT_NOTES = OUT_DIR / "MoonRockStack_wall_group_meeting_20260622_stage_success_trends_notes.md"
OUT_SUMMARY = OUT_DIR / "MoonRockStack_wall_group_meeting_20260622_stage_success_trends.summary.json"


IMG_SEED602_FRONT = ROOT / (
    "experiments/moon_rock_stack/batch_runs/"
    "20260621_course3net_upperheuristic_4to5_moon_candidates5_seed602_v1/"
    "captures_960x720_20260622/"
    "01_single_face_wall_4course_v1_failure_statics_wall_moon_trial_00/"
    "wall_front_rgb.png"
)
IMG_SEED602_TOP_DEPTH = ROOT / (
    "experiments/moon_rock_stack/batch_runs/"
    "20260621_course3net_upperheuristic_4to5_moon_candidates5_seed602_v1/"
    "captures_960x720_20260622/"
    "01_single_face_wall_4course_v1_failure_statics_wall_moon_trial_00/"
    "wall_top_depth.png"
)
IMG_SEED604_FRONT = ROOT / (
    "experiments/moon_rock_stack/batch_runs/"
    "20260622_course3net_upperheuristic_4to5_moon_candidates5_seed604_w050_v1/"
    "captures_960x720_20260622/"
    "00_single_face_wall_4course_v1_failure_statics_wall_moon_trial_00/"
    "wall_front_rgb.png"
)
IMG_SEED604_TOP_DEPTH = ROOT / (
    "experiments/moon_rock_stack/batch_runs/"
    "20260622_course3net_upperheuristic_4to5_moon_candidates5_seed604_w050_v1/"
    "captures_960x720_20260622/"
    "00_single_face_wall_4course_v1_failure_statics_wall_moon_trial_00/"
    "wall_top_depth.png"
)


FONT = "Microsoft YaHei"
BG = RGBColor(248, 250, 252)
INK = RGBColor(17, 24, 39)
MUTED = RGBColor(75, 85, 99)
SUBTLE = RGBColor(229, 231, 235)
BLUE = RGBColor(37, 99, 235)
GREEN = RGBColor(22, 163, 74)
ORANGE = RGBColor(234, 88, 12)
RED = RGBColor(220, 38, 38)
PURPLE = RGBColor(124, 58, 237)
TEAL = RGBColor(13, 148, 136)
AMBER = RGBColor(245, 158, 11)
WHITE = RGBColor(255, 255, 255)


def ensure_dirs() -> None:
    ASSET_DIR.mkdir(parents=True, exist_ok=True)


def set_bg(slide, color=BG) -> None:
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = color


def add_text(
    slide,
    text: str,
    x: float,
    y: float,
    w: float,
    h: float,
    *,
    size: int = 18,
    bold: bool = False,
    color: RGBColor = INK,
    align: PP_ALIGN = PP_ALIGN.LEFT,
    valign: MSO_ANCHOR = MSO_ANCHOR.TOP,
    font: str = FONT,
    line_spacing: float | None = 1.08,
) -> object:
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.vertical_anchor = valign
    p = tf.paragraphs[0]
    p.alignment = align
    if line_spacing is not None:
        p.line_spacing = line_spacing
    run = p.add_run()
    run.text = text
    run.font.name = font
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    return tb


def add_title(slide, title: str, subtitle: str | None = None) -> None:
    add_text(slide, title, 0.55, 0.25, 12.05, 0.45, size=25, bold=True, color=INK)
    if subtitle:
        add_text(slide, subtitle, 0.58, 0.75, 12.0, 0.28, size=9, color=MUTED)
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.55), Inches(1.08), Inches(12.2), Inches(0.018))
    line.fill.solid()
    line.fill.fore_color.rgb = SUBTLE
    line.line.fill.background()


def add_section_label(slide, label: str, x: float, y: float, color: RGBColor = BLUE) -> None:
    shp = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(1.4), Inches(0.28))
    shp.fill.solid()
    shp.fill.fore_color.rgb = color
    shp.line.fill.background()
    tf = shp.text_frame
    tf.clear()
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    r.text = label
    r.font.name = FONT
    r.font.size = Pt(8)
    r.font.bold = True
    r.font.color.rgb = WHITE


def add_card(slide, x: float, y: float, w: float, h: float, title: str, body: str, color: RGBColor = BLUE) -> None:
    shp = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    shp.fill.solid()
    shp.fill.fore_color.rgb = WHITE
    shp.line.color.rgb = SUBTLE
    shp.line.width = Pt(1)
    accent = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(y), Inches(0.06), Inches(h))
    accent.fill.solid()
    accent.fill.fore_color.rgb = color
    accent.line.fill.background()
    add_text(slide, title, x + 0.18, y + 0.13, w - 0.3, 0.23, size=12, bold=True, color=color)
    add_text(slide, body, x + 0.18, y + 0.45, w - 0.32, h - 0.52, size=9, color=INK, line_spacing=1.03)


def add_metric(slide, x: float, y: float, w: float, h: float, label: str, value: str, note: str, color: RGBColor) -> None:
    shp = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    shp.fill.solid()
    shp.fill.fore_color.rgb = WHITE
    shp.line.color.rgb = SUBTLE
    shp.line.width = Pt(1)
    add_text(slide, value, x + 0.13, y + 0.12, w - 0.26, 0.33, size=22, bold=True, color=color, align=PP_ALIGN.CENTER)
    add_text(slide, label, x + 0.13, y + 0.55, w - 0.26, 0.25, size=9, bold=True, color=INK, align=PP_ALIGN.CENTER)
    add_text(slide, note, x + 0.13, y + 0.85, w - 0.26, 0.3, size=7, color=MUTED, align=PP_ALIGN.CENTER)


def add_bullets(
    slide,
    bullets: Iterable[str],
    x: float,
    y: float,
    w: float,
    h: float,
    *,
    size: int = 11,
    color: RGBColor = INK,
    bullet_color: RGBColor = BLUE,
    line_spacing: float = 1.05,
) -> object:
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.clear()
    tf.word_wrap = True
    first = True
    for item in bullets:
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.line_spacing = line_spacing
        p.space_after = Pt(2)
        p.level = 0
        r1 = p.add_run()
        r1.text = "■ "
        r1.font.name = FONT
        r1.font.size = Pt(size)
        r1.font.color.rgb = bullet_color
        r2 = p.add_run()
        r2.text = item
        r2.font.name = FONT
        r2.font.size = Pt(size)
        r2.font.color.rgb = color
    return tb


def add_table(
    slide,
    x: float,
    y: float,
    w: float,
    h: float,
    headers: list[str],
    rows: list[list[str]],
    *,
    col_widths: list[float] | None = None,
    font_size: int = 8,
    header_color: RGBColor = BLUE,
) -> object:
    table_shape = slide.shapes.add_table(len(rows) + 1, len(headers), Inches(x), Inches(y), Inches(w), Inches(h))
    table = table_shape.table
    if col_widths:
        for i, cw in enumerate(col_widths):
            table.columns[i].width = Inches(cw)
    for j, header in enumerate(headers):
        cell = table.cell(0, j)
        cell.fill.solid()
        cell.fill.fore_color.rgb = header_color
        cell.text = header
        for p in cell.text_frame.paragraphs:
            p.alignment = PP_ALIGN.CENTER
            for r in p.runs:
                r.font.name = FONT
                r.font.size = Pt(font_size)
                r.font.bold = True
                r.font.color.rgb = WHITE
    for i, row in enumerate(rows, start=1):
        for j, text in enumerate(row):
            cell = table.cell(i, j)
            cell.fill.solid()
            cell.fill.fore_color.rgb = RGBColor(255, 255, 255) if i % 2 else RGBColor(243, 244, 246)
            cell.text = text
            for p in cell.text_frame.paragraphs:
                p.alignment = PP_ALIGN.LEFT if j == 0 else PP_ALIGN.CENTER
                for r in p.runs:
                    r.font.name = FONT
                    r.font.size = Pt(font_size)
                    r.font.color.rgb = INK
    return table_shape


def add_image(slide, path: Path, x: float, y: float, w: float, h: float, label: str | None = None) -> None:
    if path.exists():
        slide.shapes.add_picture(str(path), Inches(x), Inches(y), width=Inches(w), height=Inches(h))
        if label:
            add_text(slide, label, x, y + h + 0.05, w, 0.22, size=7, color=MUTED, align=PP_ALIGN.CENTER)
    else:
        shp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
        shp.fill.solid()
        shp.fill.fore_color.rgb = RGBColor(241, 245, 249)
        shp.line.color.rgb = SUBTLE
        add_text(slide, f"缺失图像\n{path}", x + 0.15, y + 0.2, w - 0.3, h - 0.4, size=8, color=RED)


def create_charts() -> dict[str, Path]:
    plt.rcParams["font.sans-serif"] = ["Microsoft YaHei", "SimHei", "Arial Unicode MS", "DejaVu Sans"]
    plt.rcParams["axes.unicode_minus"] = False
    charts: dict[str, Path] = {}

    roles = ["base", "middle", "cap"]
    succ = [277, 234, 96]
    fail = [7, 299, 81]
    skip = [2, 40, 25]
    fig, ax = plt.subplots(figsize=(7.0, 3.2), dpi=160)
    ax.bar(roles, succ, label="success", color="#16a34a")
    ax.bar(roles, fail, bottom=succ, label="failure", color="#dc2626")
    ax.bar(roles, skip, bottom=[a + b for a, b in zip(succ, fail)], label="skipped", color="#f59e0b")
    ax.set_title("v17 placement samples by role")
    ax.set_ylabel("placements")
    ax.legend(ncol=3, loc="upper right", frameon=False)
    ax.grid(axis="y", alpha=0.25)
    fig.tight_layout()
    charts["role_stack"] = ASSET_DIR / "role_success_failure_v17.png"
    fig.savefig(charts["role_stack"])
    plt.close(fig)

    courses = ["c0", "c1", "c2", "c3", "c4", "c5", "c6", "c7"]
    c_succ = [277, 104, 92, 100, 32, 0, 0, 0]
    c_fail = [7, 128, 145, 89, 13, 2, 0, 2]
    c_skip = [2, 13, 18, 25, 9, 0, 0, 0]
    fig, ax = plt.subplots(figsize=(7.0, 3.2), dpi=160)
    ax.bar(courses, c_succ, label="success", color="#16a34a")
    ax.bar(courses, c_fail, bottom=c_succ, label="failure", color="#dc2626")
    ax.bar(courses, c_skip, bottom=[a + b for a, b in zip(c_succ, c_fail)], label="skipped", color="#f59e0b")
    ax.set_title("v17 placement samples by course")
    ax.set_ylabel("placements")
    ax.legend(ncol=3, loc="upper right", frameon=False)
    ax.grid(axis="y", alpha=0.25)
    fig.tight_layout()
    charts["course_stack"] = ASSET_DIR / "course_success_failure_v17.png"
    fig.savefig(charts["course_stack"])
    plt.close(fig)

    metrics = {
        "v15 top1": 0.440,
        "v15 top3": 0.924,
        "v16 top1": 0.396,
        "v16 top3": 0.911,
        "critic top1": 0.362,
        "critic top3": 0.877,
        "risk acc": 0.790,
        "risk F1": 0.882,
    }
    fig, ax = plt.subplots(figsize=(7.0, 3.0), dpi=160)
    colors = ["#2563eb", "#60a5fa", "#2563eb", "#60a5fa", "#7c3aed", "#a78bfa", "#0d9488", "#5eead4"]
    ax.bar(list(metrics.keys()), list(metrics.values()), color=colors)
    ax.set_ylim(0, 1.0)
    ax.set_title("Offline neural metrics are diagnostic, not yet closed-loop sufficient")
    ax.grid(axis="y", alpha=0.25)
    for tick in ax.get_xticklabels():
        tick.set_rotation(28)
        tick.set_ha("right")
    fig.tight_layout()
    charts["network_metrics"] = ASSET_DIR / "network_metrics_v15_v16.png"
    fig.savefig(charts["network_metrics"])
    plt.close(fig)

    runs = [
        ("4L-601", 0.2616, 0.2264, 0.0451, 0.0270),
        ("4L-602", 0.3611, 0.1158, 0.4136, 0.1446),
        ("4L-604", 0.2589, 0.2209, 0.0954, 0.0108),
        ("5L-603", 0.3623, 0.2385, 0.0473, 1.6987),
        ("5L-605", 0.2211, 0.7923, 0.0126, 0.0221),
        ("5L-606", 0.2764, 0.9562, 0.0031, 0.0337),
    ]
    labels = [r[0] for r in runs]
    height = [r[1] for r in runs]
    rmse = [r[2] for r in runs]
    drift = [r[3] for r in runs]
    vel = [r[4] for r in runs]
    fig, ax = plt.subplots(figsize=(7.0, 3.2), dpi=160)
    ax.plot(labels, height, marker="o", label="height m", color="#16a34a", linewidth=2)
    ax.plot(labels, rmse, marker="o", label="RMSE m", color="#dc2626", linewidth=2)
    ax.plot(labels, drift, marker="o", label="drift m", color="#f59e0b", linewidth=2)
    ax.plot(labels, vel, marker="o", label="velocity", color="#7c3aed", linewidth=2)
    ax.set_title("Height, wall-line error, drift and velocity trade off")
    ax.legend(ncol=4, loc="upper left", frameon=False, fontsize=7)
    ax.grid(axis="y", alpha=0.25)
    fig.tight_layout()
    charts["run_tradeoffs"] = ASSET_DIR / "run_tradeoffs_20260622.png"
    fig.savefig(charts["run_tradeoffs"])
    plt.close(fig)

    return charts


def add_flow_diagram(slide) -> None:
    y = 1.48
    boxes = [
        ("岩石几何先验", "主面/棱角\n包围盒/粗糙度", BLUE),
        ("堆叠区观测", "support map\n俯视深度/占据", TEAL),
        ("候选位姿生成", "有限 yaw/roll\n接触/重心约束", ORANGE),
        ("小网络打分", "StoneSlotNet\nSupportMap/PoseRisk", PURPLE),
        ("MuJoCo 验证", "地球/月球重力\n漂移/速度/墙线", GREEN),
        ("数据回流", "正负样本\nhard negatives", RED),
    ]
    x0 = 0.65
    w = 1.82
    gap = 0.23
    for i, (title, body, color) in enumerate(boxes):
        x = x0 + i * (w + gap)
        shp = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(1.22))
        shp.fill.solid()
        shp.fill.fore_color.rgb = WHITE
        shp.line.color.rgb = color
        shp.line.width = Pt(1.3)
        add_text(slide, title, x + 0.08, y + 0.13, w - 0.16, 0.23, size=10, bold=True, color=color, align=PP_ALIGN.CENTER)
        add_text(slide, body, x + 0.08, y + 0.46, w - 0.16, 0.55, size=8, color=INK, align=PP_ALIGN.CENTER)
        if i < len(boxes) - 1:
            arr = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, Inches(x + w - 0.03), Inches(y + 0.46), Inches(gap + 0.06), Inches(0.28))
            arr.fill.solid()
            arr.fill.fore_color.rgb = RGBColor(148, 163, 184)
            arr.line.fill.background()
    add_bullets(
        slide,
        [
            "运行时输入只允许使用放置前可获得的信息：岩石几何、当前墙体观测、目标槽位、候选位姿、重力/层级。",
            "仿真后的成功率、漂移、速度、墙线误差只作为监督标签和实验日志，不能作为真实部署输入。",
            "当前策略不是端到端大模型，而是多个小网络逐步替代启发式中最贵、最不稳定的部分。",
        ],
        0.78,
        3.15,
        11.7,
        1.55,
        size=12,
        bullet_color=BLUE,
    )


def add_timeline(slide) -> None:
    stages = [
        ("0", "随机/候选试放", "能收集失败，但真实任务代价过高", RED),
        ("1", "干叠几何启发式", "主面、重心、支承重叠、墙线误差", ORANGE),
        ("2", "前三层小网络", "StoneSlotNet + support-map CNN + PoseRisk", BLUE),
        ("3", "4-5层混合闭环", "上层保留启发式和 MuJoCo 验证", TEAL),
        ("4", "下一代状态评价器", "预测未来支承、墙线和动力学风险", GREEN),
    ]
    x0, y0 = 0.75, 1.55
    w, h, gap = 2.25, 1.0, 0.2
    for i, (idx, title, body, color) in enumerate(stages):
        x = x0 + i * (w + gap)
        circle = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(x + 0.02), Inches(y0 - 0.42), Inches(0.38), Inches(0.38))
        circle.fill.solid()
        circle.fill.fore_color.rgb = color
        circle.line.fill.background()
        add_text(slide, idx, x + 0.02, y0 - 0.36, 0.38, 0.2, size=8, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y0), Inches(w), Inches(h))
        box.fill.solid()
        box.fill.fore_color.rgb = WHITE
        box.line.color.rgb = color
        box.line.width = Pt(1.2)
        add_text(slide, title, x + 0.12, y0 + 0.12, w - 0.24, 0.22, size=10, bold=True, color=color)
        add_text(slide, body, x + 0.12, y0 + 0.42, w - 0.24, 0.42, size=7, color=INK)
        if i < len(stages) - 1:
            line = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, Inches(x + w - 0.02), Inches(y0 + 0.38), Inches(gap + 0.05), Inches(0.25))
            line.fill.solid()
            line.fill.fore_color.rgb = RGBColor(148, 163, 184)
            line.line.fill.background()


def add_notes(slide_notes: list[tuple[str, list[str]]]) -> None:
    lines: list[str] = [
        "# MoonRockStack 2026-06-22 阶段性成功与发展趋势讲稿",
        "",
        "说明：PPT 文件保留 2026-06-19 旧版 15 页作为背景，新增页从第 16 页开始。",
        "本讲稿用于组会汇报，包含页面上放不下的实验目的、判断依据和口头解释。",
        "",
    ]
    for idx, (title, bullets) in enumerate(slide_notes, start=16):
        lines.append(f"## 第 {idx} 页：{title}")
        lines.append("")
        for bullet in bullets:
            lines.append(f"- {bullet}")
        lines.append("")
    OUT_NOTES.write_text("\n".join(lines), encoding="utf-8")


def build_deck() -> None:
    ensure_dirs()
    charts = create_charts()
    prs = Presentation(str(BASE_PPTX))
    blank = prs.slide_layouts[6]
    notes: list[tuple[str, list[str]]] = []

    def new_slide(title: str, subtitle: str | None = None):
        slide = prs.slides.add_slide(blank)
        set_bg(slide)
        add_title(slide, title, subtitle)
        return slide

    # 16
    slide = new_slide("新增章节：阶段性成功与发展趋势", "2026-06-22 暂停实验后的组会汇报补充；旧版 15 页保留作为前置背景")
    add_section_label(slide, "新增", 0.62, 1.27, BLUE)
    add_text(slide, "这一阶段的核心结论", 0.72, 1.65, 5.8, 0.45, size=24, bold=True, color=INK)
    add_bullets(
        slide,
        [
            "已经从“随机试错堆石头”推进到“前三层网络辅助 + 高层启发式/MuJoCo 验证 + 正负样本回流”的数据飞轮。",
            "目前最重要的成果不是 strict success=1 的完整墙，而是获得了可解释的 4 层阶段性结构案例和大量 hard negative。",
            "下一步研究重心应从盲目冲高转向提高 3-4 层闭环成功率，让网络逐步替代启发式搜索的候选排序部分。",
        ],
        0.75,
        2.2,
        5.9,
        2.45,
        size=13,
        bullet_color=BLUE,
    )
    add_image(slide, IMG_SEED602_FRONT, 7.05, 1.35, 5.35, 3.95, "阶段性结构案例 seed602：有 4 层竖向结构，但 drift 未过严格阈值")
    notes.append(
        (
            "新增章节：阶段性成功与发展趋势",
            [
                "向老师先说明：这不是把失败包装成成功。本阶段 strict wall success 还没有达到稳定水平，但已经出现可复用的中间成果。",
                "阶段性成果包括：有效图像记录、明确失败模式、可训练的数据集、以及能支撑前三层的成熟小网络组合。",
                "汇报口径要强调科学实验的严谨性：失败样本不是无用尝试，而是后续网络学习墙线、支承和动力学风险的监督来源。",
            ],
        )
    )

    # 17
    slide = new_slide("失败样本到成功样本：当前数据账本", "v17 数据集已经从 run 级失败中提取 placement/candidate 级正负样本")
    add_metric(slide, 0.7, 1.35, 1.75, 1.15, "run examples", "41", "完整实验记录", BLUE)
    add_metric(slide, 2.7, 1.35, 1.75, 1.15, "placements", "1061", "逐石放置样本", TEAL)
    add_metric(slide, 4.7, 1.35, 1.75, 1.15, "candidate poses", "34836", "候选位姿样本", PURPLE)
    add_metric(slide, 6.7, 1.35, 1.75, 1.15, "Moon positives", "592", "月面正样本", GREEN)
    add_metric(slide, 8.7, 1.35, 1.75, 1.15, "Moon negatives", "378", "月面失败样本", RED)
    add_metric(slide, 10.7, 1.35, 1.75, 1.15, "skipped", "67", "无可行候选", AMBER)
    add_image(slide, charts["role_stack"], 0.82, 2.85, 5.45, 2.65, "按角色统计：base 不是瓶颈，middle/cap 是主要学习对象")
    add_image(slide, charts["course_stack"], 6.8, 2.85, 5.45, 2.65, "按层统计：从第 2-3 层开始失败显著上升")
    add_text(
        slide,
        "解释：run 级 strict success 仍少，但一次失败 run 内仍含有大量局部成功/局部失败放置；这些样本正是训练候选排序网络的主要来源。",
        0.82,
        5.75,
        11.45,
        0.48,
        size=11,
        color=INK,
    )
    notes.append(
        (
            "失败样本到成功样本：当前数据账本",
            [
                "回答老师可能会问的第一个问题：多少失败才产生正样本？目前 v17 有 36 个 run 目录、41 条 run example、1061 条 placement、34836 条候选位姿。",
                "月面 placement 层面已经有 592 条 success、378 条 failure、67 条 skipped。也就是说，虽然完整墙严格成功仍不稳定，但局部放置正样本已经足够支持小网络迭代。",
                "base 的失败很少：286 个 base 样本中 success 277、failure 7。middle 和 cap 才是高价值失败来源：middle failure 299、cap failure 81。",
            ],
        )
    )

    # 18
    slide = new_slide("阶段性成功案例：4 层结构已经出现，但还不是严格成功", "seed602 是当前最适合汇报的阶段性正案例：结构感明显，但 drift 暴露了下一阶段问题")
    add_image(slide, IMG_SEED602_FRONT, 0.7, 1.35, 5.45, 4.05, "正视 RGB：形成竖向 4 层局部结构")
    add_image(slide, IMG_SEED602_TOP_DEPTH, 6.75, 1.35, 5.45, 4.05, "俯视 depth：墙线/柱状偏移可见")
    add_table(
        slide,
        0.9,
        5.62,
        11.0,
        0.65,
        ["run", "stable/failure", "RMSE", "height", "drift", "velocity", "判断"],
        [["seed602 4-course", "20/4", "0.1158 m", "0.3611 m", "0.4136 m", "0.1446", "阶段性结构，严格失败"]],
        col_widths=[1.9, 1.3, 1.2, 1.2, 1.2, 1.2, 2.8],
        font_size=8,
        header_color=GREEN,
    )
    notes.append(
        (
            "阶段性成功案例：4 层结构已经出现，但还不是严格成功",
            [
                "这一页要非常严谨：seed602 不能说成完整成功。它在视觉上有 4 层结构，stable/failure 是 20/4，高度 0.3611 m，RMSE 0.1158 m，说明墙线较好。",
                "但是 drift 0.4136 m，说明最终结构整体横向漂移明显，strict success 仍失败。",
                "该案例的价值是证明：前三层网络 + 高层启发式能够把系统推到 4 层结构附近，下一步应解决 drift 和全局墙线保持。",
            ],
        )
    )

    # 19
    slide = new_slide("典型负案例：数值稳定不等于墙体成功", "seed604 证明只压 drift/velocity 会得到“稳定但错误的形状”")
    add_image(slide, IMG_SEED604_FRONT, 0.7, 1.35, 5.45, 4.05, "正视 RGB：局部堆散，缺少连续单面墙")
    add_image(slide, IMG_SEED604_TOP_DEPTH, 6.75, 1.35, 5.45, 4.05, "俯视 depth：支承点散开，墙线目标没有被保持")
    add_table(
        slide,
        0.9,
        5.62,
        11.0,
        0.65,
        ["run", "stable/failure", "RMSE", "height", "drift", "velocity", "判断"],
        [["seed604 4-course", "17/7", "0.2209 m", "0.2589 m", "0.0954 m", "0.0108", "动力学稳，但不是墙"]],
        col_widths=[1.9, 1.3, 1.2, 1.2, 1.2, 1.2, 2.8],
        font_size=8,
        header_color=RED,
    )
    notes.append(
        (
            "典型负案例：数值稳定不等于墙体成功",
            [
                "seed604 使用更高 PoseRisk 权重后，drift 和 velocity 明显降低，但图像显示结构是局部石堆，不是单面墙。",
                "这直接回答了哪些规则没用：单独增加稳定性惩罚、单独延长 settling、单独增加候选数量，都不能保证墙线形状。",
                "这个失败样本应作为 hard negative 进入训练：标签不仅惩罚动力学风险，也要惩罚墙线 RMSE、横向 spread 和下一层可支承性。",
            ],
        )
    )

    # 20
    slide = new_slide("哪些规则有用，哪些规则暂时无效", "规则不是最终答案，但它们定义了训练数据、损失函数和安全约束")
    add_table(
        slide,
        0.72,
        1.35,
        12.0,
        4.8,
        ["类别", "规则/方法", "观察到的效果", "后续处理"],
        [
            ["有用", "棱角化多面体、拒绝尖刺/光滑圆边", "更符合石头几何先验，避免非物理支承", "保留为生成器硬约束"],
            ["有用", "主面/支承面、重心投影、支承重叠", "能显著减少明显会滑落的候选", "作为网络输入和损失项"],
            ["有用", "前三层网络裁剪 top-3，4-5 层保留启发式", "降低早期搜索成本，避免高层误剪可行候选", "当前默认策略"],
            ["有用", "正视 RGB + 俯视 depth + RMSE/drift/velocity 联合评价", "阻止把石堆误判为墙", "作为日志和报告固定模板"],
            ["无效/不足", "随机放置或随机候选搜索", "能产生数据，但真实任务开销不可接受", "只保留为探索基线"],
            ["无效/不足", "只增加 settling 或候选数量", "会得到稳定但错误的墙线", "改为 hard negative 监督"],
            ["无效/不足", "只压 drift/velocity", "低漂移不等于结构正确", "加入墙线和未来支承目标"],
            ["无效/不足", "新 v15/v16 网络直接接管闭环", "top-1 仍低，容易剪掉好候选", "先做离线诊断和辅助排序"],
        ],
        col_widths=[1.2, 3.3, 3.6, 3.0],
        font_size=7,
        header_color=BLUE,
    )
    notes.append(
        (
            "哪些规则有用，哪些规则暂时无效",
            [
                "有效规则的共同点：它们都来自放置前可获得的几何信息和当前堆叠观测，能迁移到真实系统。",
                "暂时无效规则的共同点：它们只优化局部稳定或增加搜索量，没有表达最终结构目标，所以容易得到稳定但错误的石堆。",
                "向老师说明：当前启发式不是要永久保留，而是作为网络训练前的先验、标签构造和安全约束。",
            ],
        )
    )

    # 21
    slide = new_slide("神经网络是在什么阶段接入的", "最初不是端到端控制，而是多个小网络替代候选选择中的高成本环节")
    add_timeline(slide)
    add_table(
        slide,
        0.85,
        3.25,
        11.55,
        2.35,
        ["网络", "最初接入位置", "输入", "输出", "目的"],
        [
            ["StoneSlotNet", "选石头/匹配槽位", "石头几何 + 槽位/层级", "石头-槽位分数", "减少不合适石头进入候选池"],
            ["support-map CNN", "候选位姿排序", "局部支承高度/占据/目标/footprint", "候选 pose rank", "把 top-k 候选留给 MuJoCo"],
            ["PoseRiskNet", "风险惩罚", "候选几何 + 位姿指标 + 支承指标", "风险分数", "降低滑移、漂移、残余速度风险"],
        ],
        col_widths=[1.65, 2.1, 3.1, 1.7, 2.75],
        font_size=7,
        header_color=PURPLE,
    )
    notes.append(
        (
            "神经网络是在什么阶段接入的",
            [
                "最初接入方式是保守的模块化小网络，不是直接输入图片输出最终放置动作。",
                "目的有三个：第一，减少昂贵的候选搜索；第二，避免随机试错在真实机器人上不可用；第三，让失败经验转化为可学习的候选排序信号。",
                "当前最成熟的闭环是前三层网络化，高层仍用启发式和短仿真验证，这是为了避免 top-1 不够强的新网络误删高层可行候选。",
            ],
        )
    )

    # 22
    slide = new_slide("神经网络输入/输出如何演变", "从单石头几何拟合，演变到“石头 + 当前墙体观测 + 候选位姿 + 未来支承”的联合评价")
    add_flow_diagram(slide)
    notes.append(
        (
            "神经网络输入/输出如何演变",
            [
                "早期只看石头外观/几何来选石头，容易变成数据拟合：某类石头历史上成功，所以以后也选它。",
                "现在的设计明确避免把测试后验作为运行时输入。运行时只输入放置前能知道的信息：岩石几何、当前墙体深度/support map、目标槽位、候选位姿、重力和层级。",
                "输出也从简单的成功/失败，逐步演变为多目标：当前是否稳定、是否贴合墙线、是否给下一层留下支承、是否会在月面低重力下产生漂移或残余速度。",
            ],
        )
    )

    # 23
    slide = new_slide("新网络当前能做什么，不能做什么", "v15/v16 的 top-3 说明候选池里有答案；top-1 不够说明还不能完全替代启发式")
    add_image(slide, charts["network_metrics"], 0.8, 1.35, 5.8, 2.65, "离线指标：top-3 高，但 top-1 仍不足")
    add_table(
        slide,
        7.0,
        1.35,
        5.35,
        2.65,
        ["模型", "关键指标", "闭环判断"],
        [
            ["v15 support-map CNN", "top1 0.440 / top3 0.924", "可辅助分析，不替换 v4"],
            ["v16 support-map CNN", "top1 0.396 / top3 0.911", "hard negative 后泛化未提升"],
            ["v16 WallStateCritic", "top1 0.362 / top3 0.877", "离线诊断可用，闭环裁剪不足"],
            ["v16 PoseRiskNet", "acc 0.790 / F1 0.882", "风险标签偏正，仍不如旧 v5 稳"],
        ],
        col_widths=[1.9, 1.85, 1.6],
        font_size=7,
        header_color=PURPLE,
    )
    add_bullets(
        slide,
        [
            "能做：减少前三层候选、生成更稳定的中高层训练分布、帮助识别 hard negative。",
            "不能做：目前还不能直接输出高层最终位姿，也不能完全取代 MuJoCo 验证和墙线约束。",
            "发展趋势：从候选局部 ranker 升级为 wall-state/future-support critic，再逐步接管更高层的候选排序。",
        ],
        0.9,
        4.35,
        11.2,
        1.35,
        size=12,
        bullet_color=PURPLE,
    )
    notes.append(
        (
            "新网络当前能做什么，不能做什么",
            [
                "top-3 高的意义：好候选通常在网络前几名里，说明网络学到了一部分几何和支承规律。",
                "top-1 低的意义：如果让网络直接决定唯一动作，仍然会经常选错，所以真实部署还不能跳过启发式和物理验证。",
                "当前科学价值在于：网络已经能参与数据飞轮，但下一步要提高 top-1 和全局结构评价，而不是盲目换大模型。",
            ],
        )
    )

    # 24
    slide = new_slide("第二层、第三层、第四层的难题有什么不同", "层数越高，问题越从局部接触稳定转向全局误差累积和未来支承")
    add_table(
        slide,
        0.72,
        1.28,
        12.0,
        4.75,
        ["层级", "当前统计", "主要难题", "与上一层不同点", "当前处理"],
        [
            ["第 1 层 base", "277/286 success", "不是支承问题，而是基础 footprint 和墙线起点", "地面平整，主要决定后续结构边界", "几何筛选 + 规则放置"],
            ["第 2 层 course1", "104/245 success, 128 failure", "第一次真正面对不平支承；yaw/roll 和重心投影敏感", "从地面支承变成石-石支承", "网络 top-3 + PoseRisk"],
            ["第 3 层 course2", "92/255 success, 145 failure", "需要跨接、互锁，还要给第 4 层留下可支承面", "不只是当前稳定，还要维护未来支承", "前三层网络辅助核心区域"],
            ["第 4 层 course3", "100/214 success, 89 failure", "下层误差已累积；局部稳定可能破坏整体墙线", "全局墙线/drift 开始主导 strict fail", "高层回退启发式 + MuJoCo"],
            ["第 5 层及 cap", "course4: 32/54 success", "支承选项少，易形成稳定但错误的形状", "需要闭合顶部并不推散墙体", "记录 hard negative，训练未来支承 critic"],
        ],
        col_widths=[1.25, 2.0, 3.1, 2.8, 2.2],
        font_size=7,
        header_color=TEAL,
    )
    add_text(
        slide,
        "共同问题：都需要足够支承面积、重心落在支承多边形内、低残余速度。不同问题：层数越高，墙线形状和下一层可支承性越重要。",
        0.9,
        6.12,
        11.2,
        0.38,
        size=10,
        color=INK,
    )
    notes.append(
        (
            "第二层、第三层、第四层的难题有什么不同",
            [
                "第二层主要是从地面支承过渡到石-石支承，问题集中在支承面积、姿态和重心投影。",
                "第三层开始出现误差传播和未来支承问题：一个石头当下能稳住，不代表它会给第四层留下可用平台。",
                "第四层更难的原因是全局目标变强：局部稳定候选很多，但能同时保持墙线、减少 drift、保留上层支承的候选很少。",
            ],
        )
    )

    # 25
    slide = new_slide("4-5 层实验揭示的主要矛盾", "高度、墙线、漂移、速度之间不是单调兼容，需要多目标训练")
    add_image(slide, charts["run_tradeoffs"], 0.75, 1.28, 6.0, 2.75, "多指标 trade-off：一个指标好不能代表墙成功")
    add_table(
        slide,
        7.05,
        1.28,
        5.45,
        2.75,
        ["现象", "例子", "含义"],
        [
            ["高度不错但速度高", "5L seed603: height 0.3623, vel 1.6987", "月面低重力残余运动仍危险"],
            ["漂移低但墙线差", "5L seed605: drift 0.0126, RMSE 0.7923", "稳定地堆歪不是成功"],
            ["候选更多仍失败", "5L seed606: candidates=8, RMSE 0.9562", "搜索量不是结构目标"],
            ["4 层结构但 drift 大", "4L seed602: height 0.3611, drift 0.4136", "需要全局姿态保持"],
        ],
        col_widths=[1.6, 2.05, 1.8],
        font_size=7,
        header_color=ORANGE,
    )
    add_bullets(
        slide,
        [
            "如果只奖励高度，系统会生成高但会动/会倒的结构。",
            "如果只惩罚 drift/velocity，系统会选择低矮、横向散开的稳定石堆。",
            "下一代损失函数必须同时包含墙线 RMSE、支承连续性、未来层可支承面积和动力学风险。",
        ],
        0.9,
        4.35,
        11.3,
        1.3,
        size=12,
        bullet_color=ORANGE,
    )
    notes.append(
        (
            "4-5 层实验揭示的主要矛盾",
            [
                "4-5 层不是简单把 3 层成功重复一次。层数越高，目标之间的冲突越明显。",
                "高度、墙线、漂移、速度必须联合评价。任一单指标被过度优化，都会产生看似稳定但结构错误的结果。",
                "因此后续网络目标不应只是分类 success/failure，而应该预测一组结构指标，或者用多任务损失学习综合评分。",
            ],
        )
    )

    # 26
    slide = new_slide("为什么不继续盲目冲 10 层", "先把 3-4 层成功率做上去，科学价值高于偶然高层案例")
    add_card(slide, 0.75, 1.35, 3.75, 1.45, "当前问题", "4-5 层 strict success 仍不稳定；继续冲 8-10 层会产生偶然案例，但难以解释和复现。", RED)
    add_card(slide, 4.85, 1.35, 3.75, 1.45, "阶段目标", "把 3 层闭环成功率提高到 60%-80%，再让网络参与 4-5 层候选排序。", BLUE)
    add_card(slide, 8.95, 1.35, 3.75, 1.45, "研究价值", "形成可复用的数据飞轮：失败变成负样本，成功变成正样本，启发式逐步退出主决策。", GREEN)
    add_table(
        slide,
        0.9,
        3.25,
        11.25,
        2.35,
        ["阶段", "主要动作", "成功判据", "产出"],
        [
            ["阶段 A", "稳定 2-3 层网络辅助堆叠", "placement success 和 wall-line 指标提升", "干净正负样本"],
            ["阶段 B", "前 3 层网络 + 第 4 层启发式/MuJoCo", "4 层结构案例增多，失败可解释", "第 4 层 hard negatives"],
            ["阶段 C", "训练 future-support / wall-state critic", "top-1 替代启发式部分候选排序", "4-5 层闭环提升"],
            ["阶段 D", "再冲 6-10 层", "不是偶然高，而是可复现实验成功率", "高墙/路标结构策略"],
        ],
        col_widths=[1.25, 3.5, 3.6, 2.9],
        font_size=8,
        header_color=BLUE,
    )
    notes.append(
        (
            "为什么不继续盲目冲 10 层",
            [
                "本阶段研究判断是：偶然冲高没有科学价值。这里用 PPT 解释为什么暂停冲高、转向提高 3-4 层成功率。",
                "真正目标是让网络学到可迁移的几何和堆叠规则，而不是记住某批石头或某些随机种子。",
                "只有当 3-4 层成功率和可解释性提高后，再做 5 层以上高墙才有说服力。",
            ],
        )
    )

    # 27
    slide = new_slide("数据飞轮和异步调度设计", "本端 2080Ti 负责主训练，远端 1080Ti/CPU 负责轻量采样；当前已暂停新实验，先完成汇报")
    add_flow_diagram(slide)
    add_table(
        slide,
        0.85,
        4.88,
        11.55,
        1.05,
        ["资源", "角色", "当前优先任务", "为什么这样分配"],
        [
            ["本端 RTX 2080Ti + 64G", "主训练/模型评估", "support-map、PoseRisk、WallStateCritic 小网络", "训练吞吐更好，便于快速迭代"],
            ["远端 GTX 1080Ti", "轻量采样/辅助训练", "低风险 run、候选日志汇总", "不阻塞本端主训练，可补充数据"],
            ["MuJoCo 仿真", "物理验证", "地球/月球重力、短仿真稳定性", "确保网络标签来自物理闭环"],
        ],
        col_widths=[2.1, 2.1, 3.5, 3.6],
        font_size=8,
        header_color=TEAL,
    )
    notes.append(
        (
            "数据飞轮和异步调度设计",
            [
                "这页说明实验组织方式：不是单次脚本，而是采样、训练、评估、回流的循环。",
                "本端 2080Ti 适合做主训练，远端 1080Ti 适合轻量采样和补充数据。为了响应当前暂停指令，本页只讲设计和已形成的调度思想，不启动新实验。",
                "强调所有数据都保留，包括失败、skipped、图像和候选日志，因为负样本是后续学习的核心。",
            ],
        )
    )

    # 28
    slide = new_slide("后续网络趋势：从候选 ranker 到结构策略", "模型要学习的是“在当前石堆上，这块石头放哪里能形成下一层可用结构”")
    add_table(
        slide,
        0.72,
        1.28,
        12.0,
        4.8,
        ["阶段", "输入", "输出", "训练标签", "预期作用"],
        [
            ["当前 A: StoneSlotNet", "石头几何 + 目标槽位", "石头选择分数", "放置成功/失败、角色适配", "减少不合适石头"],
            ["当前 B: support-map CNN", "局部深度/support map + candidate footprint", "候选位姿 rank", "候选是否接近启发式最优", "替代部分位姿搜索"],
            ["当前 C: PoseRiskNet", "候选指标 + 支承指标 + 几何", "风险惩罚", "drift/velocity/failure 标签", "减少动力学失败"],
            ["下一步 D: WallStateCritic", "当前墙体观测 + 放置后预测状态", "结构质量分数", "RMSE/spread/visible courses", "避免石堆化"],
            ["下一步 E: FutureSupportNet", "墙顶局部形状 + 候选石头", "下一层可支承性", "下一层 success/skip 标签", "让第三层服务第四层"],
            ["远期 F: 端到端策略", "点云/深度 + 石头库 + 目标结构", "石头 ID + 6D 位姿", "闭环成功和结构奖励", "减少启发式主导"],
        ],
        col_widths=[1.9, 3.0, 2.0, 2.8, 2.3],
        font_size=7,
        header_color=GREEN,
    )
    add_text(
        slide,
        "关键约束：历史成功率不能作为运行时输入；它只能作为训练后的评估统计。模型必须依赖几何先验和当前观测来泛化到新石头。",
        0.9,
        6.1,
        11.2,
        0.38,
        size=10,
        color=INK,
    )
    notes.append(
        (
            "后续网络趋势：从候选 ranker 到结构策略",
            [
                "这一页直接回答神经网络的发展趋势：不是马上训练一个大模型，而是从小网络逐步扩展输入和输出。",
                "输入会从单石头几何，扩展到当前墙体的 depth/support map，再扩展到未来支承能力。",
                "输出会从石头选择和候选排序，演变为结构质量、风险和未来可支承性的联合预测，最后才考虑端到端策略。",
            ],
        )
    )

    # 29
    slide = new_slide("给老师汇报时的阶段性结论", "这一阶段的成功在于明确了问题边界、数据来源和下一步可验证路线")
    add_card(slide, 0.78, 1.35, 5.55, 1.2, "1. 阶段性成功", "已经能产生 4 层局部墙/柱状结构案例；虽然 strict fail，但足以证明当前策略能进入高层结构状态。", GREEN)
    add_card(slide, 6.75, 1.35, 5.55, 1.2, "2. 失败的科学价值", "v17 已有 34836 个候选位姿和 1061 个逐石样本，middle/cap 失败为后续学习提供 hard negatives。", RED)
    add_card(slide, 0.78, 2.85, 5.55, 1.2, "3. 有效规则", "棱角多面体、主面支承、重心/支承重叠、前三层网络化、正视 RGB + 俯视 depth 联合评价。", BLUE)
    add_card(slide, 6.75, 2.85, 5.55, 1.2, "4. 暂时无效规则", "随机放置、单纯增加候选、单纯延长 settling、单独压 drift/velocity，都不能保证真正墙体结构。", ORANGE)
    add_card(slide, 0.78, 4.35, 5.55, 1.2, "5. 网络路线", "从 StoneSlotNet/support-map/PoseRisk 小网络开始，逐步发展到 WallStateCritic 和 FutureSupportNet。", PURPLE)
    add_card(slide, 6.75, 4.35, 5.55, 1.2, "6. 下一步", "先提高 3-4 层成功率，再推进 4-5 层网络闭环；等成功率可复现后再冲 6-10 层。", TEAL)
    notes.append(
        (
            "给老师汇报时的阶段性结论",
            [
                "最后用这一页收束：本阶段不是最终工程成功，而是科学实验路径变清楚了。",
                "最有说服力的点是：我们已经知道哪些规则有效、哪些规则会误导系统，也知道网络应该在哪些环节逐步接管。",
                "下一步汇报可以围绕一个可检验假设：加入 future-support 和 wall-state 监督后，3-4 层闭环成功率是否显著提升。",
            ],
        )
    )

    prs.save(str(OUT_PPTX))
    add_notes(notes)

    summary = {
        "base_pptx": str(BASE_PPTX),
        "output_pptx": str(OUT_PPTX),
        "notes_md": str(OUT_NOTES),
        "asset_dir": str(ASSET_DIR),
        "old_slide_count": 15,
        "new_slide_count": len(prs.slides),
        "added_slide_count": len(notes),
        "key_images": {
            "seed602_front": str(IMG_SEED602_FRONT),
            "seed602_top_depth": str(IMG_SEED602_TOP_DEPTH),
            "seed604_front": str(IMG_SEED604_FRONT),
            "seed604_top_depth": str(IMG_SEED604_TOP_DEPTH),
        },
        "charts": {k: str(v) for k, v in charts.items()},
    }
    OUT_SUMMARY.write_text(json.dumps(summary, ensure_ascii=False, indent=2), encoding="utf-8")
    print(json.dumps(summary, ensure_ascii=False, indent=2))


if __name__ == "__main__":
    build_deck()
